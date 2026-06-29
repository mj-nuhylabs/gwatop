"""PPTX / DOCX → 평문 텍스트 추출.

`pdf_text.extract_text_from_pdf_bytes` 와 같은 역할의 자매 모듈.
추출 결과는 PDF와 동일하게 `File.extracted_text` 에 저장되어 이후
강의계획서 파싱(`syllabus_parser`)·주차 분류(`classification`)의 입력으로 쓰인다.
따라서 출력은 "사람이 읽는 순서의 평문" 이면 충분하고, 별도 전처리는
`pdf_text.clean_syllabus_text` 가 다운스트림에서 그대로 담당한다.

- PPTX (`python-pptx`): 슬라이드 순서대로 도형 텍스트 + 표 셀 + 발표자 노트.
- DOCX (`python-docx`): 문단 + 표 셀(문서 본문 순서는 보존하지 않고 문단 → 표 순으로 단순 결합).
"""

from __future__ import annotations

import logging
from io import BytesIO

logger = logging.getLogger(__name__)

# PDF 와 동일한 페이지/슬라이드 구분자 — 다운스트림 전처리가 동일하게 동작하도록 맞춘다.
SLIDE_SEPARATOR = "\n\n\f\n\n"


def extract_text_from_pptx_bytes(data: bytes) -> str:
    """PPTX 바이트 → 평문. 슬라이드 단위로 \\f 구분자로 join."""
    from pptx import Presentation  # 지연 import — 워커 부팅 비용/선택 의존성 격리

    prs = Presentation(BytesIO(data))
    slides_text: list[str] = []

    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_iter_shape_text(shape))
        # 발표자 노트 — 강의 보충 설명이 들어 있는 경우가 많아 포함한다.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes is not None and notes.text and notes.text.strip():
                parts.append(notes.text.strip())
        if parts:
            slides_text.append("\n".join(parts))

    return SLIDE_SEPARATOR.join(slides_text)


def _iter_shape_text(shape) -> list[str]:
    """도형 하나에서 텍스트 라인들을 뽑는다. 표/그룹 도형은 재귀 처리."""
    out: list[str] = []

    # 표 도형
    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(c for c in cells if c)
            if line:
                out.append(line)
        return out

    # 그룹 도형 — 내부 도형 재귀
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub in shape.shapes:
            out.extend(_iter_shape_text(sub))
        return out

    # 일반 텍스트 프레임
    if shape.has_text_frame:
        text = shape.text_frame.text
        if text and text.strip():
            out.append(text.strip())

    return out


def extract_text_from_docx_bytes(data: bytes) -> str:
    """DOCX 바이트 → 평문. 문단 텍스트 + 표 셀 텍스트를 결합."""
    from docx import Document  # 지연 import

    doc = Document(BytesIO(data))
    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text and para.text.strip():
            parts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(c for c in cells if c)
            if line:
                parts.append(line)

    return "\n".join(parts)
